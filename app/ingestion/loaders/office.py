import logfire

import docx as python_docx
from pptx import Presentation


def _parse_docx(file_path: str) -> str:
    with logfire.span("📄 DOCX Parsing", filename=file_path):
        try:
            doc = python_docx.Document(file_path)
            full_text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
            if not full_text.strip():
                logfire.warning(f"⚠️ Unstructured returned empty text for {file_path}")
            else:
                logfire.info(f"✅ Successfully parsed {len(full_text)} characters")
            return full_text
        except Exception as e:
            logfire.error(f"❌ DOCX Parse Failed: {e}")
            raise e
    


def _parse_pptx(file_path: str) -> str:
    with logfire.span("📄 PPTX Parsing", filename=file_path):
        try:
            prs = Presentation(file_path)
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        texts.append(shape.text.strip())
            return "\n".join(texts)
        except Exception as e:
            logfire.error(f"❌ PPTX Parse Failed: {e}")
            raise e

def parse_office(file_path: str):
    with logfire.span("📄 Office Document Parsing", filename=file_path):
        try:
            if file_path.endswith(".docx"):
                return _parse_docx(file_path)
            elif file_path.endswith(".pptx"):
                return _parse_pptx(file_path)
            else:
                logfire.warning(f"⚠️ Unsupported file type: {file_path}")
                return ""
        except Exception as e:
            logfire.error(f"❌ Office Parse Failed: {e}")
            raise e

